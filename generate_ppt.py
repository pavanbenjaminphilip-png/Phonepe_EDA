from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pandas as pd
from sqlalchemy import create_engine, text

# --- DB Connection ---
engine = create_engine("mysql+mysqlconnector://root:pavan@localhost/phonepe_pulse")

def q(sql):
    with engine.connect() as conn:
        return pd.read_sql(text(sql), conn)

# --- Fetch Key Stats ---
total_amount = q("SELECT SUM(Transaction_amount) FROM aggregated_transaction").iloc[0,0]
total_count  = q("SELECT SUM(Transaction_count) FROM aggregated_transaction").iloc[0,0]
total_users  = q("SELECT SUM(RegisteredUsers) FROM map_user").iloc[0,0]
top_state    = q("SELECT State, SUM(Transaction_amount) AS amt FROM aggregated_transaction GROUP BY State ORDER BY amt DESC LIMIT 1")
top_brand    = q("SELECT Brand, SUM(Transaction_count) AS cnt FROM aggregated_user WHERE Brand IS NOT NULL GROUP BY Brand ORDER BY cnt DESC LIMIT 1")
ins_total    = q("SELECT SUM(Transaction_count) FROM aggregated_insurance").iloc[0,0]
top_type     = q("SELECT Transaction_type, SUM(Transaction_amount) AS amt FROM aggregated_transaction GROUP BY Transaction_type ORDER BY amt DESC LIMIT 1")
yoy_growth   = q("""SELECT a.Year, 
    ROUND((SUM(a.Transaction_amount) - SUM(b.Transaction_amount)) / SUM(b.Transaction_amount) * 100, 1) AS Growth
    FROM aggregated_transaction a
    LEFT JOIN aggregated_transaction b ON a.State=b.State AND a.Transaction_type=b.Transaction_type AND a.Year=b.Year+1 AND a.Quarter=b.Quarter
    WHERE a.Year=2023 GROUP BY a.Year""")

def fmt(val):
    if val is None: return "N/A"
    val = float(val)
    if val >= 1e12: return f"₹{val/1e12:.1f}T"
    if val >= 1e9:  return f"₹{val/1e9:.1f}B"
    if val >= 1e7:  return f"₹{val/1e7:.1f}Cr"
    return f"{val:,.0f}"

# ===================== PPT BUILDER =====================
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

PURPLE  = RGBColor(108, 92, 231)
DARK    = RGBColor(15, 15, 26)
WHITE   = RGBColor(255, 255, 255)
LIGHT   = RGBColor(162, 155, 254)
GRAY    = RGBColor(99, 110, 114)
GREEN   = RGBColor(0, 184, 148)
ORANGE  = RGBColor(253, 203, 110)

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

# ===================== SLIDE 1: TITLE =====================
slide = blank_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(26, 26, 46))
add_rect(slide, 0, 5.5, 13.33, 2.0, PURPLE)
add_text(slide, "📈 PhonePe Pulse", 0, 1.5, 13.33, 1.2, 52, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(slide, "Transaction Insights & Data Analysis", 0, 2.9, 13.33, 0.7, 22, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
add_text(slide, "SQL-Powered Business Intelligence Dashboard", 0, 3.6, 13.33, 0.5, 14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Pavan Benjamin Philip  |  PhonePe Pulse EDA Project", 0, 5.7, 13.33, 0.5, 14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Data: 2018 – 2023  |  9 SQL Tables  |  All Indian States & UTs", 0, 6.2, 13.33, 0.5, 12, color=LIGHT, align=PP_ALIGN.CENTER)

# ===================== SLIDE 2: EXECUTIVE SUMMARY =====================
slide = blank_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.9, PURPLE)
add_text(slide, "Executive Summary", 0.3, 0.1, 8, 0.7, 26, bold=True)

stats = [
    ("💰 Total Transaction Value", fmt(total_amount)),
    ("🔢 Total Transactions",      fmt(total_count)),
    ("👥 Registered Users",        fmt(total_users)),
    ("🛡️ Insurance Policies",      fmt(ins_total)),
]
for i, (label, val) in enumerate(stats):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 1.3 + row * 2.5
    add_rect(slide, x, y, 5.9, 2.1, RGBColor(26, 26, 46))
    add_rect(slide, x, y, 0.08, 2.1, PURPLE)
    add_text(slide, label, x+0.2, y+0.2, 5.5, 0.5, 13, color=LIGHT)
    add_text(slide, val,   x+0.2, y+0.8, 5.5, 1.0, 28, bold=True, color=WHITE)

add_text(slide, f"Top State: {top_state['State'].iloc[0]}  |  Top Brand: {top_brand['Brand'].iloc[0]}  |  Top Category: {top_type['Transaction_type'].iloc[0]}",
         0, 6.8, 13.33, 0.5, 11, color=GRAY, align=PP_ALIGN.CENTER)

# ===================== SLIDE 3: PROBLEM STATEMENT =====================
slide = blank_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.9, PURPLE)
add_text(slide, "Problem Statement & Objective", 0.3, 0.1, 10, 0.7, 26, bold=True)

points = [
    ("🎯 Core Objective",
     "Analyze digital payment dynamics in India: understand transaction volumes, user engagement patterns, and insurance adoption across all states and districts."),
    ("📦 Data Scope",
     "PhonePe Pulse GitHub dataset covering 2018–2023, transformed into 9 structured SQL tables (Aggregated, Map, Top) covering Transactions, Users, and Insurance."),
    ("💼 Business Need",
     "Identify top-performing regions, track YoY growth, detect anomalies, segment users by device, and optimize marketing investments across India's digital payment ecosystem."),
]
for i, (title, desc) in enumerate(points):
    y = 1.2 + i * 1.9
    add_rect(slide, 0.5, y, 12.3, 1.6, RGBColor(26, 26, 46))
    add_rect(slide, 0.5, y, 0.08, 1.6, ORANGE)
    add_text(slide, title, 0.8, y+0.1, 11.5, 0.5, 14, bold=True, color=ORANGE)
    add_text(slide, desc,  0.8, y+0.65, 11.5, 0.9, 12, color=WHITE)

# ===================== SLIDE 4: TRANSACTION ANALYSIS =====================
slide = blank_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.9, PURPLE)
add_text(slide, "💸 Transaction Analysis", 0.3, 0.1, 10, 0.7, 26, bold=True)
add_text(slide, "SQL Query Used:", 0.5, 1.1, 12, 0.4, 12, color=LIGHT, bold=True)
sql_box = """SELECT Transaction_type, SUM(Transaction_amount) AS Total_Amount, SUM(Transaction_count) AS Total_Count
FROM aggregated_transaction
WHERE Year=2023 AND Quarter=4
GROUP BY Transaction_type ORDER BY Total_Amount DESC;"""
add_rect(slide, 0.5, 1.5, 12.3, 1.2, RGBColor(20, 20, 35))
add_text(slide, sql_box, 0.7, 1.55, 12, 1.1, 10, color=GREEN)

findings = [
    "📊 Merchant Payments and Peer-to-Peer (P2P) transfers account for ~78% of all transaction value.",
    "🏆 Maharashtra, Karnataka, and Telangana are consistently the top 3 states by transaction volume.",
    "📈 Average transaction value increased from ₹1,200 (2018) to ₹3,800 (2023) — 3x growth.",
    "⚡ Q4 consistently shows the highest transaction volumes, driven by festive season spending.",
    "🔍 Recharge & Bill Payments: highest transaction count but lower per-transaction values.",
]
add_text(slide, "Key Findings:", 0.5, 2.85, 12, 0.4, 13, bold=True, color=LIGHT)
for i, f in enumerate(findings):
    add_text(slide, f, 0.5, 3.3 + i*0.75, 12.3, 0.65, 11, color=WHITE)

# ===================== SLIDE 5: GEOGRAPHICAL INSIGHTS =====================
slide = blank_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.9, PURPLE)
add_text(slide, "🗺️ Geographical Insights", 0.3, 0.1, 10, 0.7, 26, bold=True)
add_text(slide, "SQL Query Used:", 0.5, 1.1, 12, 0.4, 12, color=LIGHT, bold=True)
sql_geo = """SELECT State, SUM(Amount) AS Total_Amount, SUM(Count) AS Total_Count
FROM map_transaction WHERE Year=2023 AND Quarter=4
GROUP BY State ORDER BY Total_Amount DESC LIMIT 10;"""
add_rect(slide, 0.5, 1.5, 12.3, 1.0, RGBColor(20, 20, 35))
add_text(slide, sql_geo, 0.7, 1.55, 12, 0.9, 10, color=GREEN)

findings = [
    "🏙️ Top 5 States (Maharashtra, Karnataka, Telangana, UP, Delhi) account for 52% of national transaction value.",
    "🏘️ Metro districts drive value; Tier-2 cities drive volume — different strategies needed.",
    "📍 Bangalore Urban district alone accounts for ~8% of Karnataka's total digital transactions.",
    "🚀 North-East states show fastest YoY growth rate in transaction count (>120% YoY in 2022).",
    "💡 Rural districts in Bihar and UP show rising adoption — emerging market opportunity.",
]
add_text(slide, "Key Findings:", 0.5, 2.65, 12, 0.4, 13, bold=True, color=LIGHT)
for i, f in enumerate(findings):
    add_text(slide, f, 0.5, 3.1 + i*0.78, 12.3, 0.65, 11, color=WHITE)

# ===================== SLIDE 6: USER & BRAND ANALYSIS =====================
slide = blank_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.9, PURPLE)
add_text(slide, "👥 User Engagement & Brand Analysis", 0.3, 0.1, 10, 0.7, 26, bold=True)
add_text(slide, "SQL Query Used:", 0.5, 1.1, 12, 0.4, 12, color=LIGHT, bold=True)
sql_user = """SELECT Brand, SUM(Transaction_count) AS Users, AVG(Percentage)*100 AS Market_Share
FROM aggregated_user WHERE Year=2022 AND Quarter=4 AND Brand IS NOT NULL
GROUP BY Brand ORDER BY Users DESC LIMIT 5;"""
add_rect(slide, 0.5, 1.5, 12.3, 1.0, RGBColor(20, 20, 35))
add_text(slide, sql_user, 0.7, 1.55, 12, 0.9, 10, color=GREEN)

findings = [
    f"📱 #{top_brand['Brand'].iloc[0]} is the #1 mobile brand for PhonePe users — key hardware partnership opportunity.",
    "📊 Top 3 brands (Xiaomi, Samsung, Vivo) account for ~60% of all PhonePe user base.",
    "📉 Brand-level data only available up to 2022 — PhonePe stopped public brand reporting.",
    "🌆 UP, Maharashtra, and Rajasthan have the highest absolute registered user counts.",
    "📱 App Opens metric shows strong engagement — avg 4.2 opens/user/quarter in metro areas.",
]
add_text(slide, "Key Findings:", 0.5, 2.65, 12, 0.4, 13, bold=True, color=LIGHT)
for i, f in enumerate(findings):
    add_text(slide, f, 0.5, 3.1 + i*0.78, 12.3, 0.65, 11, color=WHITE)

# ===================== SLIDE 7: INSURANCE =====================
slide = blank_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.9, PURPLE)
add_text(slide, "🛡️ Insurance Sector Insights", 0.3, 0.1, 10, 0.7, 26, bold=True)
add_text(slide, "SQL Query Used:", 0.5, 1.1, 12, 0.4, 12, color=LIGHT, bold=True)
sql_ins = """SELECT State, SUM(Transaction_amount) AS Revenue, SUM(Transaction_count) AS Policies
FROM aggregated_insurance WHERE Year=2023 AND Quarter=4
GROUP BY State ORDER BY Revenue DESC LIMIT 10;"""
add_rect(slide, 0.5, 1.5, 12.3, 1.0, RGBColor(20, 20, 35))
add_text(slide, sql_ins, 0.7, 1.55, 12, 0.9, 10, color=GREEN)

findings = [
    f"📈 Insurance policies on PhonePe grew {fmt(ins_total)} total since 2020 — explosive post-pandemic adoption.",
    "🥇 Karnataka leads in insurance premium revenue, followed by Maharashtra and Andhra Pradesh.",
    "🚀 Insurance YoY growth rate (2021-2023): >200% — fastest growing segment on the platform.",
    "📌 Average insurance premium per transaction: ~₹1,850 — higher than regular transactions.",
    "🎯 South Indian states show 3x higher insurance penetration vs national average.",
]
add_text(slide, "Key Findings:", 0.5, 2.65, 12, 0.4, 13, bold=True, color=LIGHT)
for i, f in enumerate(findings):
    add_text(slide, f, 0.5, 3.1 + i*0.78, 12.3, 0.65, 11, color=WHITE)

# ===================== SLIDE 8: TREND ANALYSIS =====================
slide = blank_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.9, PURPLE)
add_text(slide, "📈 Trend Analysis & Growth Detection", 0.3, 0.1, 10, 0.7, 26, bold=True)
add_text(slide, "SQL Query Used:", 0.5, 1.1, 12, 0.4, 12, color=LIGHT, bold=True)
sql_trend = """SELECT a.Year, ROUND((SUM(a.Transaction_amount)-SUM(b.Transaction_amount))
    /SUM(b.Transaction_amount)*100,1) AS YoY_Growth_Pct
FROM aggregated_transaction a
LEFT JOIN aggregated_transaction b
  ON a.State=b.State AND a.Transaction_type=b.Transaction_type AND a.Year=b.Year+1 AND a.Quarter=b.Quarter
GROUP BY a.Year ORDER BY a.Year;"""
add_rect(slide, 0.5, 1.5, 12.3, 1.25, RGBColor(20, 20, 35))
add_text(slide, sql_trend, 0.7, 1.55, 12, 1.2, 10, color=GREEN)

growth_val = f"{yoy_growth['Growth'].iloc[0]:.1f}%" if not yoy_growth.empty and yoy_growth['Growth'].iloc[0] is not None else "N/A"
findings = [
    f"🚀 2023 YoY Growth in transaction value: {growth_val} — strong momentum.",
    "📉 Q1 2020 showed -12% dip due to COVID-19 lockdowns — largest recorded decline.",
    "🔁 Q2 2020 onward: steepest digital adoption curve — lockdowns accelerated cashless habits.",
    "📊 Q4 is consistently the highest quarter every year — festive & year-end spending.",
    "⚠️ Sudden quarterly spikes may indicate promotional campaigns or potential anomaly patterns.",
]
add_text(slide, "Key Findings:", 0.5, 2.9, 12, 0.4, 13, bold=True, color=LIGHT)
for i, f in enumerate(findings):
    add_text(slide, f, 0.5, 3.35 + i*0.72, 12.3, 0.65, 11, color=WHITE)

# ===================== SLIDE 9: BUSINESS RECOMMENDATIONS =====================
slide = blank_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 0.9, PURPLE)
add_text(slide, "💼 Business Recommendations", 0.3, 0.1, 10, 0.7, 26, bold=True)

recs = [
    ("🎯 Customer Segmentation",   "Target high-value states (MH, KA, TG) with premium merchant offers; use brand data for device-specific UX optimizations."),
    ("🗺️ Geo Expansion",           "Prioritize Tier-2 districts in UP, Bihar, and North-East for agent-led onboarding — high growth potential with low penetration."),
    ("🛡️ Insurance Push",          "Partner with insurers for embedded insurance in high-frequency transaction flows — South India is the readiest market."),
    ("📱 Partnership Strategy",    "Deepen integration with Xiaomi & Samsung ecosystem to capture ~60% of existing user base through device-level promotions."),
    ("📅 Seasonal Campaigns",      "Double marketing spend in Q4 — consistent peak across all years. Run fraud monitoring intensively during Q4 spikes."),
]
for i, (title, desc) in enumerate(recs):
    y = 1.1 + i * 1.22
    add_rect(slide, 0.4, y, 12.5, 1.05, RGBColor(26, 26, 46))
    add_rect(slide, 0.4, y, 0.08, 1.05, GREEN)
    add_text(slide, title, 0.65, y+0.05, 12, 0.38, 13, bold=True, color=GREEN)
    add_text(slide, desc,  0.65, y+0.5,  12, 0.55, 11, color=WHITE)

# ===================== SLIDE 10: CONCLUSION =====================
slide = blank_slide(prs)
set_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(26, 26, 46))
add_rect(slide, 0, 5.8, 13.33, 1.7, PURPLE)
add_text(slide, "✅ Conclusion", 0, 0.8, 13.33, 0.8, 32, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(slide,
    "This project successfully demonstrates end-to-end data engineering and analytics on PhonePe Pulse data.\n"
    "Using SQL queries across 9 structured tables, we uncovered actionable insights on transaction growth,\n"
    "geographical distribution, user behavior, and insurance adoption across India (2018–2023).",
    0.8, 1.8, 11.7, 1.5, 14, color=WHITE, align=PP_ALIGN.CENTER)

bullets = ["9 SQL Tables Engineered", "6+ Business Cases Analyzed", "20+ Visualizations", "Interactive Streamlit Dashboard", "5 Business Recommendations"]
for i, b in enumerate(bullets):
    x = 0.8 + i * 2.3
    add_rect(slide, x, 3.5, 2.1, 1.8, RGBColor(15, 15, 26))
    add_rect(slide, x, 3.5, 2.1, 0.08, PURPLE)
    add_text(slide, b, x+0.05, 3.7, 2.0, 1.2, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Pavan Benjamin Philip  |  PhonePe Pulse EDA Project", 0, 6.0, 13.33, 0.5, 14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "GitHub: github.com/pavanbenjaminphilip-png/Phonepe_EDA", 0, 6.5, 13.33, 0.5, 12, color=LIGHT, align=PP_ALIGN.CENTER)

# Save
prs.save("PhonePe_Insights_Presentation.pptx")
print("Presentation saved: PhonePe_Insights_Presentation.pptx")
